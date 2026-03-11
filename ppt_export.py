"""
ppt_export.py – Paste screenshot images into PowerPoint slides.

Requirements:
    pip install python-pptx
"""

from __future__ import annotations

import io
from collections import defaultdict
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple, Union


class PPTExportError(Exception):
    """Raised when a PowerPoint export operation fails."""


class PPTExporter:
    """Insert images into PowerPoint slides using python-pptx."""

    _EMU_PER_INCH = 914400

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def paste_image(
        self,
        pptx_path: str,
        slide_number: int,                    # 1-based
        image_source: Union[str, io.BytesIO], # file path OR in-memory bytes
        left: float,                          # inches
        top: float,
        width: float,
        height: float,
    ) -> None:
        """
        Insert an image on slide *slide_number* of *pptx_path*.

        *image_source* can be:
          - a file path (str) – read from disk
          - an io.BytesIO   – pasted directly from memory (clipboard capture)
        """
        self._check_exists(pptx_path, "PowerPoint file")
        if isinstance(image_source, str):
            self._check_exists(image_source, "Screenshot")

        prs = self._open(pptx_path)
        n = len(prs.slides)
        if not (1 <= slide_number <= n):
            raise PPTExportError(
                f"Slide {slide_number} does not exist.\n"
                f"The file has {n} slide{'s' if n != 1 else ''}."
            )
        try:
            from pptx.util import Inches  # noqa: PLC0415
            slide = prs.slides[slide_number - 1]
            slide.shapes.add_picture(
                image_source,
                Inches(left), Inches(top),
                Inches(width), Inches(height),
            )
        except PPTExportError:
            raise
        except Exception as exc:
            raise PPTExportError(f"Cannot add picture to slide: {exc}") from exc

        self._save(prs, pptx_path)

    def paste_batch(
        self, jobs: List[dict], log=None
    ) -> List[Tuple[dict, str | None]]:
        """
        Paste multiple images, grouping by pptx_path so each file is
        opened and saved only once.

        Each job dict must have keys:
            pptx_path, slide_number, image_source (str path OR io.BytesIO),
            left, top, width, height, entry (the Entry object for reporting)

        Returns list of (job, error_message_or_None).
        """
        try:
            from pptx.util import Inches  # noqa: PLC0415
        except ImportError:
            raise PPTExportError(
                "python-pptx is not installed.\nFix: pip install python-pptx"
            )

        _log = log or (lambda msg, level="info": None)

        # Group by file so each PPTX is opened/saved only once
        by_file: Dict[str, List[dict]] = defaultdict(list)
        for job in jobs:
            by_file[job["pptx_path"]].append(job)

        results = []
        for pptx_path, file_jobs in by_file.items():
            _log(f"Opening  {Path(pptx_path).name}", "info")
            try:
                self._check_exists(pptx_path, "PowerPoint file")
                prs = self._open(pptx_path)
                n = len(prs.slides)
            except PPTExportError as exc:
                _log(f"Cannot open: {exc}", "err")
                for j in file_jobs:
                    results.append((j, str(exc)))
                continue

            for j in file_jobs:
                sn = j["slide_number"]
                entry = j.get("entry")
                label = entry.name if entry else pptx_path
                if not (1 <= sn <= n):
                    msg = f"Slide {sn} does not exist (file has {n})."
                    _log(f"  ✗ {label}: {msg}", "err")
                    results.append((j, msg))
                    continue
                try:
                    img_src = j["image_source"]
                    if isinstance(img_src, str):
                        self._check_exists(img_src, "Screenshot")
                    elif hasattr(img_src, "seek"):
                        img_src.seek(0)   # rewind BytesIO before use
                    slide = prs.slides[sn - 1]
                    slide.shapes.add_picture(
                        img_src,
                        Inches(j["left"]), Inches(j["top"]),
                        Inches(j["width"]), Inches(j["height"]),
                    )
                    _log(f"  Pasted  '{label}'  →  slide {sn}", "ok")
                    results.append((j, None))
                except Exception as exc:
                    _log(f"  ✗ {label}: {exc}", "err")
                    results.append((j, str(exc)))

            try:
                _log(f"Saving  {Path(pptx_path).name}…", "info")
                self._save(prs, pptx_path)
                _log(f"Saved  {Path(pptx_path).name}", "ok")
            except PPTExportError as exc:
                _log(f"Save failed: {exc}", "err")
                for j in file_jobs:
                    for idx, (rj, rerr) in enumerate(results):
                        if rj is j and rerr is None:
                            results[idx] = (j, str(exc))

        return results

    @staticmethod
    def get_slide_info(pptx_path: str) -> Tuple[int, float, float]:
        """
        Return (slide_count, width_inches, height_inches).
        Returns (0, 10.0, 7.5) on any error.
        """
        _EMU = 914400
        try:
            from pptx import Presentation  # noqa: PLC0415
            prs = Presentation(pptx_path)
            w = round(prs.slide_width / _EMU, 2)
            h = round(prs.slide_height / _EMU, 2)
            return len(prs.slides), w, h
        except Exception:
            return 0, 10.0, 7.5

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _check_exists(path: str, label: str) -> None:
        if not Path(path).exists():
            raise PPTExportError(f"{label} not found:\n{path}")

    @staticmethod
    def _open(pptx_path: str):
        try:
            from pptx import Presentation  # noqa: PLC0415
            return Presentation(pptx_path)
        except ImportError:
            raise PPTExportError(
                "python-pptx is not installed.\nFix: pip install python-pptx"
            )
        except Exception as exc:
            raise PPTExportError(f"Cannot open PowerPoint file:\n{exc}") from exc

    @staticmethod
    def _save(prs, pptx_path: str) -> None:
        try:
            prs.save(pptx_path)
        except Exception as exc:
            raise PPTExportError(f"Cannot save PowerPoint file:\n{exc}") from exc
